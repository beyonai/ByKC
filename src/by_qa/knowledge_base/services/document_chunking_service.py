"""Server-side document chunking and embedding generation."""

import csv
import io
from dataclasses import dataclass, field
from pathlib import PurePosixPath

import httpx
from langchain_text_splitters import (
    MarkdownHeaderTextSplitter,
    RecursiveCharacterTextSplitter,
)

from by_qa.knowledge_base.api.schemas import KnowledgeItemChunkPayload
from by_qa.knowledge_base.services.errors import KnowledgeBaseConfigurationError

SUPPORTED_EXTENSIONS = {
    ".md",
    ".markdown",
    ".txt",
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
}


@dataclass
class DocumentChunkingService:
    """Extract text from documents, split into chunks, and generate embeddings."""

    embedding_base_url: str
    embedding_api_key: str
    embedding_model_name: str
    embedding_dimension: int
    chunk_size: int = 512
    chunk_overlap: int = 64
    embedding_timeout: float = 60.0
    _headers_to_split_on: list[tuple[str, str]] = field(
        default_factory=lambda: [
            ("#", "h1"),
            ("##", "h2"),
            ("###", "h3"),
            ("####", "h4"),
        ]
    )

    FILE_TYPE_TO_EXT = {
        "pdf": ".pdf",
        "docx": ".docx",
        "pptx": ".pptx",
        "xlsx": ".xlsx",
    }

    def extract_text_from_file(self, file_bytes: bytes, file_type: str) -> str:
        """Extract text from a file given its type label (e.g. 'pdf', 'docx')."""
        ext = self.FILE_TYPE_TO_EXT.get(file_type)
        if ext is None:
            supported = ", ".join(sorted(self.FILE_TYPE_TO_EXT))
            raise ValueError(
                f"unsupported file type: {file_type}. Supported types: {supported}"
            )
        return self._extract_text(file_bytes, ext)

    def chunk_and_embed(
        self, file_bytes: bytes, *, filename: str
    ) -> list[KnowledgeItemChunkPayload]:
        """Extract text, split into chunks, and generate embeddings."""
        ext = PurePosixPath(filename).suffix.lower()
        text = self._extract_text(file_bytes, ext)
        if not text or not text.strip():
            raise ValueError("document produced no text after extraction")

        chunks = self._split_text(text, ext)
        if not chunks:
            raise ValueError("document produced no chunks after splitting")

        for chunk in chunks:
            start_line, end_line = self._compute_line_range(text, chunk["chunk_text"])
            chunk["start_line"] = start_line
            chunk["end_line"] = end_line

        embeddings = self._batch_embed([c["chunk_text"] for c in chunks])
        return [
            KnowledgeItemChunkPayload(
                chunk_no=c["chunk_no"],
                start_line=c["start_line"],
                end_line=c["end_line"],
                chunk_text=c["chunk_text"],
                embedding=emb,
                char_start=c["char_start"],
                char_end=c["char_end"],
            )
            for c, emb in zip(chunks, embeddings)
        ]

    @staticmethod
    def _compute_line_range(full_text: str, chunk_text: str) -> tuple[int, int]:
        """Compute 1-based start_line and end_line for a chunk within the document."""
        char_start = full_text.find(chunk_text)
        if char_start < 0:
            return 1, 1

        start_line = full_text[:char_start].count("\n") + 1
        end_line = start_line + chunk_text.count("\n")
        return start_line, end_line

    # ------------------------------------------------------------------
    # Text extraction
    # ------------------------------------------------------------------

    def _extract_text(self, file_bytes: bytes, ext: str) -> str:
        if ext in (".md", ".markdown", ".txt"):
            return file_bytes.decode("utf-8")
        if ext == ".pdf":
            return self._extract_pdf(file_bytes)
        if ext == ".docx":
            return self._extract_docx(file_bytes)
        if ext == ".pptx":
            return self._extract_pptx(file_bytes)
        if ext == ".xlsx":
            return self._extract_xlsx(file_bytes)
        if ext == ".csv":
            return self._extract_csv(file_bytes)
        raise ValueError(
            f"unsupported file type: {ext}. "
            f"Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    @staticmethod
    def _extract_pdf(file_bytes: bytes) -> str:
        try:
            import fitz  # pymupdf
        except ImportError as exc:
            raise KnowledgeBaseConfigurationError(
                "pymupdf is required for PDF support: pip install pymupdf"
            ) from exc
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n\n".join(pages)

    @staticmethod
    def _extract_docx(file_bytes: bytes) -> str:
        try:
            from docx import Document
        except ImportError as exc:
            raise KnowledgeBaseConfigurationError(
                "python-docx is required for DOCX support: pip install python-docx"
            ) from exc
        doc = Document(io.BytesIO(file_bytes))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    @staticmethod
    def _extract_pptx(file_bytes: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise KnowledgeBaseConfigurationError(
                "python-pptx is required for PPTX support: pip install python-pptx"
            ) from exc
        prs = Presentation(io.BytesIO(file_bytes))
        slides: list[str] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n".join(slides)

    @staticmethod
    def _extract_xlsx(file_bytes: bytes) -> str:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise KnowledgeBaseConfigurationError(
                "openpyxl is required for Excel support: pip install openpyxl"
            ) from exc
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"## {sheet_name}\n\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    @staticmethod
    def _extract_csv(file_bytes: bytes) -> str:
        text = file_bytes.decode("utf-8")
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(row))
        return "\n".join(rows)

    # ------------------------------------------------------------------
    # Text splitting
    # ------------------------------------------------------------------

    def _split_text(self, text: str, ext: str) -> list[dict]:
        if ext in (".md", ".markdown"):
            return self._split_markdown(text)
        return self._split_plain(text)

    def _split_markdown(self, text: str) -> list[dict]:
        header_splitter = MarkdownHeaderTextSplitter(
            headers_to_split_on=self._headers_to_split_on,
            strip_headers=False,
        )
        header_docs = header_splitter.split_text(text)

        char_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
        )

        chunks: list[dict] = []
        chunk_no = 1
        for doc in header_docs:
            sub_docs = char_splitter.split_text(doc.page_content)
            for sub_text in sub_docs:
                char_start = text.find(sub_text)
                char_end = char_start + len(sub_text) if char_start >= 0 else None
                chunks.append(
                    {
                        "chunk_no": chunk_no,
                        "chunk_text": sub_text,
                        "char_start": char_start if char_start >= 0 else None,
                        "char_end": char_end,
                    }
                )
                chunk_no += 1
        return chunks

    def _split_plain(self, text: str) -> list[dict]:
        char_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
        )
        sub_docs = char_splitter.split_text(text)

        chunks: list[dict] = []
        chunk_no = 1
        for sub_text in sub_docs:
            char_start = text.find(sub_text)
            char_end = char_start + len(sub_text) if char_start >= 0 else None
            chunks.append(
                {
                    "chunk_no": chunk_no,
                    "chunk_text": sub_text,
                    "char_start": char_start if char_start >= 0 else None,
                    "char_end": char_end,
                }
            )
            chunk_no += 1
        return chunks

    # ------------------------------------------------------------------
    # Embedding
    # ------------------------------------------------------------------

    def _batch_embed(self, texts: list[str]) -> list[list[float]]:
        if not self.embedding_base_url:
            raise KnowledgeBaseConfigurationError(
                "EMBEDDING_BASE_URL is required for server-side embedding"
            )

        headers = {"Content-Type": "application/json"}
        if self.embedding_api_key:
            headers["Authorization"] = f"Bearer {self.embedding_api_key}"

        response = httpx.post(
            f"{self.embedding_base_url.rstrip('/')}/embeddings",
            headers=headers,
            json={
                "model": self.embedding_model_name,
                "input": texts,
                "dimensions": self.embedding_dimension,
            },
            timeout=self.embedding_timeout,
        )
        response.raise_for_status()
        payload = response.json()
        data = payload.get("data") or []
        if len(data) != len(texts):
            raise KnowledgeBaseConfigurationError(
                f"embedding API returned {len(data)} vectors for {len(texts)} inputs"
            )

        sorted_data = sorted(data, key=lambda d: d.get("index", 0))
        embeddings = [item["embedding"] for item in sorted_data]

        for i, emb in enumerate(embeddings):
            if len(emb) != self.embedding_dimension:
                raise KnowledgeBaseConfigurationError(
                    f"embedding dimension mismatch at chunk {i}: "
                    f"got {len(emb)}, expected {self.embedding_dimension}"
                )
        return embeddings
