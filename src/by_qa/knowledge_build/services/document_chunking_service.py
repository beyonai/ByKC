"""Server-side document chunking and embedding generation."""

import csv
import io
import re
import unicodedata
from dataclasses import dataclass, field
from pathlib import PurePosixPath

import httpx

from by_qa.core import logger
from by_qa.knowledge_build.services.heading_patterns import (
    DEFAULT_HEADING_PATTERNS,
    HeadingPattern,
)
from by_qa.knowledge_common.exceptions import KnowledgeConfigurationError
from by_qa.knowledge_common.schemas import KnowledgeItemChunkPayload

SUPPORTED_EXTENSIONS = {
    ".md",
    ".markdown",
    ".txt",
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
}


@dataclass
class _TextBlock:
    text: str
    start_char: int
    end_char: int
    start_line: int
    end_line: int
    kind: str
    level: int | None = None


@dataclass
class _HeadingMatch:
    pattern_name: str
    line_no: int
    explicit_level: int | None = None
    ordered_item_no: int | None = None


@dataclass
class DocumentChunkingService:
    """Extract text from documents, split into chunks, and generate embeddings."""

    embedding_base_url: str
    embedding_api_key: str
    embedding_model_name: str
    embedding_dimension: int
    embedding_batch_max_texts: int = 10
    chunk_size: int = 512
    chunk_overlap: int = 64
    embedding_timeout: float = 60.0
    heading_patterns: list[HeadingPattern] = field(
        default_factory=lambda: list(DEFAULT_HEADING_PATTERNS)
    )

    FILE_TYPE_TO_EXT = {
        "txt": ".txt",
        "md": ".md",
        "csv": ".csv",
        "pdf": ".pdf",
        "docx": ".docx",
        "pptx": ".pptx",
        "xlsx": ".xlsx",
    }

    def extract_text_from_file(self, file_bytes: bytes, file_type: str) -> str:
        """Extract text from a file given its type label."""
        normalized_file_type = file_type.strip().lower()
        if normalized_file_type == "markdown":
            normalized_file_type = "md"
        ext = self.FILE_TYPE_TO_EXT.get(normalized_file_type)
        if ext is None:
            supported = ", ".join(sorted(self.FILE_TYPE_TO_EXT))
            raise ValueError(
                f"unsupported file type: {file_type}. Supported types: {supported}"
            )
        return self._extract_text(file_bytes, ext)

    def chunk_and_embed(
        self, file_bytes: bytes, *, filename: str
    ) -> list[KnowledgeItemChunkPayload]:
        """Extract text, split into chunks, and generate embeddings."""
        ext = PurePosixPath(filename).suffix.lower()
        text = self._extract_text(file_bytes, ext)
        if not text or not text.strip():
            raise ValueError("document produced no text after extraction")

        chunks = self._split_text(text, ext)
        if not chunks:
            raise ValueError("document produced no chunks after splitting")
        logger.info(
            "document_chunking markdown chunking completed: filename=%s, chunk_count=%s",
            filename,
            len(chunks),
        )

        embeddings = self._batch_embed([c["chunk_text"] for c in chunks])
        logger.info(
            "document_chunking embedding completed: filename=%s, chunk_count=%s",
            filename,
            len(embeddings),
        )
        return [
            KnowledgeItemChunkPayload(
                chunk_no=c["chunk_no"],
                start_line=c["start_line"],
                end_line=c["end_line"],
                chunk_text=c["chunk_text"],
                embedding=emb,
            )
            for c, emb in zip(chunks, embeddings)
        ]

    def _extract_text(self, file_bytes: bytes, ext: str) -> str:
        if ext in (".md", ".markdown", ".txt"):
            return file_bytes.decode("utf-8")
        if ext == ".pdf":
            return self._extract_pdf(file_bytes)
        if ext == ".docx":
            return self._extract_docx(file_bytes)
        if ext == ".pptx":
            return self._extract_pptx(file_bytes)
        if ext == ".xlsx":
            return self._extract_xlsx(file_bytes)
        if ext == ".csv":
            return self._extract_csv(file_bytes)
        raise ValueError(
            f"unsupported file type: {ext}. "
            f"Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    @staticmethod
    def _extract_pdf(file_bytes: bytes) -> str:
        try:
            import fitz
        except ImportError as exc:
            raise KnowledgeConfigurationError(
                "pymupdf is required for PDF support: pip install pymupdf"
            ) from exc
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n\n".join(pages)

    @staticmethod
    def _extract_docx(file_bytes: bytes) -> str:
        try:
            from docx import Document
        except ImportError as exc:
            raise KnowledgeConfigurationError(
                "python-docx is required for DOCX support: pip install python-docx"
            ) from exc
        doc = Document(io.BytesIO(file_bytes))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))
        return "\n\n".join(parts)

    @staticmethod
    def _extract_pptx(file_bytes: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise KnowledgeConfigurationError(
                "python-pptx is required for PPTX support: pip install python-pptx"
            ) from exc
        prs = Presentation(io.BytesIO(file_bytes))
        slides: list[str] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n".join(slides)

    @staticmethod
    def _extract_xlsx(file_bytes: bytes) -> str:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise KnowledgeConfigurationError(
                "openpyxl is required for Excel support: pip install openpyxl"
            ) from exc
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"## {sheet_name}\n\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    @staticmethod
    def _extract_csv(file_bytes: bytes) -> str:
        text = file_bytes.decode("utf-8")
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(row))
        return "\n".join(rows)

    def _split_text(self, text: str, ext: str) -> list[dict]:
        blocks = self._build_blocks(text, treat_as_markdown=ext in (".md", ".markdown"))
        return self._build_chunks_from_blocks(text, blocks)

    def _build_blocks(self, text: str, *, treat_as_markdown: bool) -> list[_TextBlock]:
        lines = self._build_line_entries(text)
        noise_line_nos = self._detect_noise_line_numbers(
            lines, treat_as_markdown=treat_as_markdown
        )
        if treat_as_markdown:
            excluded_line_nos = noise_line_nos | self._detect_code_block_line_numbers(
                lines
            )
        else:
            excluded_line_nos = noise_line_nos
        heading_levels = self._infer_heading_levels(
            lines, excluded_line_nos, treat_as_markdown=treat_as_markdown
        )
        blocks: list[_TextBlock] = []
        index = 0
        while index < len(lines):
            line = lines[index]
            if line["line_no"] in noise_line_nos:
                index += 1
                continue
            stripped = line["content"].strip()
            if not stripped:
                index += 1
                continue

            heading_level = heading_levels.get(line["line_no"])
            if heading_level is not None:
                blocks.append(
                    _TextBlock(
                        text=text[line["start_char"] : line["end_char"]],
                        start_char=line["start_char"],
                        end_char=line["end_char"],
                        start_line=line["line_no"],
                        end_line=line["line_no"],
                        kind="heading",
                        level=heading_level,
                    )
                )
                index += 1
                continue

            start_index = index
            while index < len(lines):
                current = lines[index]
                if current["line_no"] in noise_line_nos:
                    index += 1
                    continue
                if not current["content"].strip():
                    break
                if heading_levels.get(current["line_no"]) is not None:
                    break
                index += 1

            kept_lines = [
                line
                for line in lines[start_index:index]
                if line["line_no"] not in noise_line_nos and line["content"].strip()
            ]
            if not kept_lines:
                index += 1
                continue
            first_line = kept_lines[0]
            last_line = kept_lines[-1]
            blocks.append(
                _TextBlock(
                    text="\n".join(line["content"] for line in kept_lines),
                    start_char=first_line["start_char"],
                    end_char=last_line["end_char"],
                    start_line=first_line["line_no"],
                    end_line=last_line["line_no"],
                    kind="paragraph",
                )
            )
        return blocks

    def _detect_noise_line_numbers(
        self, lines: list[dict], *, treat_as_markdown: bool
    ) -> set[int]:
        if treat_as_markdown:
            return set()

        stripped_counts: dict[str, int] = {}
        for line in lines:
            stripped = line["content"].strip()
            if stripped:
                stripped_counts[stripped] = stripped_counts.get(stripped, 0) + 1

        noise_line_nos: set[int] = set()
        for line in lines:
            stripped = line["content"].strip()
            if not stripped:
                continue
            if re.match(r"^\d+\s*/\s*\d+$", stripped):
                noise_line_nos.add(line["line_no"])
                continue
            if stripped_counts.get(stripped, 0) < 2:
                continue
            if stripped.endswith(".md"):
                noise_line_nos.add(line["line_no"])
                continue
            if re.match(r"^\d{4}-\d{2}-\d{2}$", stripped):
                noise_line_nos.add(line["line_no"])
                continue
        return noise_line_nos

    @staticmethod
    def _detect_code_block_line_numbers(lines: list[dict]) -> set[int]:
        """Return line numbers inside fenced code blocks (including the fence lines).

        Lines inside ``` or ~~~ fences must not be matched against heading patterns
        because code comments like '# filename.yml' would otherwise corrupt the
        heading stack and produce wrong breadcrumbs.
        """
        code_block_lines: set[int] = set()
        fence_char: str | None = None
        fence_len: int = 0

        for line in lines:
            stripped = line["content"].strip()
            if fence_char is None:
                m = re.match(r"^(`{3,}|~{3,})", stripped)
                if m:
                    fence_char = m.group(1)[0]
                    fence_len = len(m.group(1))
                    code_block_lines.add(line["line_no"])
            else:
                code_block_lines.add(line["line_no"])
                m = re.match(r"^(`{3,}|~{3,})\s*$", stripped)
                if m and m.group(1)[0] == fence_char and len(m.group(1)) >= fence_len:
                    fence_char = None
                    fence_len = 0

        return code_block_lines

    @staticmethod
    def _build_line_entries(text: str) -> list[dict]:
        lines: list[dict] = []
        start_char = 0
        for line_no, line_with_newline in enumerate(
            text.splitlines(keepends=True), start=1
        ):
            line_break_len = len(line_with_newline) - len(
                line_with_newline.rstrip("\n")
            )
            end_char = start_char + len(line_with_newline) - line_break_len
            lines.append(
                {
                    "line_no": line_no,
                    "content": line_with_newline[
                        : len(line_with_newline) - line_break_len
                    ],
                    "start_char": start_char,
                    "end_char": end_char,
                }
            )
            start_char += len(line_with_newline)

        if not lines:
            lines.append(
                {
                    "line_no": 1,
                    "content": text,
                    "start_char": 0,
                    "end_char": len(text),
                }
            )
        return lines

    def _infer_heading_levels(
        self,
        lines: list[dict],
        noise_line_nos: set[int],
        *,
        treat_as_markdown: bool,
    ) -> dict[int, int]:
        matches: list[_HeadingMatch] = []
        for line in lines:
            if line["line_no"] in noise_line_nos:
                continue
            match = self._match_heading_pattern(
                line["content"], treat_as_markdown=treat_as_markdown
            )
            if match is None:
                continue
            match.line_no = line["line_no"]
            matches.append(match)

        skipped_ordered_list_lines = self._detect_ordered_list_line_numbers(matches)

        inferred_levels: dict[str, int] = {}
        next_level = 1
        for match in matches:
            if match.line_no in skipped_ordered_list_lines:
                continue
            if match.explicit_level is not None:
                continue
            if match.pattern_name in inferred_levels:
                continue
            inferred_levels[match.pattern_name] = next_level
            next_level += 1

        heading_levels: dict[int, int] = {}
        for match in matches:
            if match.line_no in skipped_ordered_list_lines:
                continue
            if match.explicit_level is not None:
                heading_levels[match.line_no] = match.explicit_level
            else:
                heading_levels[match.line_no] = inferred_levels[match.pattern_name]
        return heading_levels

    @staticmethod
    def _detect_ordered_list_line_numbers(matches: list[_HeadingMatch]) -> set[int]:
        line_numbers: set[int] = set()
        index = 0
        while index < len(matches):
            match = matches[index]
            if match.pattern_name != "numeric_dot" or match.ordered_item_no is None:
                index += 1
                continue

            run = [match]
            cursor = index + 1
            while cursor < len(matches):
                previous = run[-1]
                current = matches[cursor]
                if (
                    current.pattern_name != "numeric_dot"
                    or current.ordered_item_no is None
                    or current.line_no != previous.line_no + 1
                    or current.ordered_item_no != previous.ordered_item_no + 1
                ):
                    break
                run.append(current)
                cursor += 1

            if len(run) >= 2:
                line_numbers.update(item.line_no for item in run)
            index = cursor

        return line_numbers

    def _match_heading_pattern(
        self, line: str, *, treat_as_markdown: bool
    ) -> _HeadingMatch | None:
        stripped = line.strip()
        normalized = unicodedata.normalize("NFKC", stripped)
        if not stripped:
            return None
        for pattern in self.heading_patterns:
            if pattern.markdown_only and not treat_as_markdown:
                continue
            if pattern.reject_if_contains_colon and (
                "：" in normalized or ":" in normalized
            ):
                continue
            if not re.match(pattern.regex, normalized):
                continue
            if not pattern.markdown_only and self._looks_like_inline_body(normalized):
                continue
            pattern_name = pattern.name
            ordered_item_no = None
            if pattern.name == "numeric_nested":
                nested_match = re.match(r"^(\d+(?:\.\d+)+)\s+(.+)$", normalized)
                if nested_match is None:
                    continue
                title_text = nested_match.group(2).strip()
                if not title_text or re.fullmatch(r"[\d.\s]+", normalized):
                    continue
                numeric_depth = nested_match.group(1).count(".") + 1
                pattern_name = f"{pattern.name}_{numeric_depth}"
            elif pattern.name == "numeric_dot":
                ordered_match = re.match(r"^(\d+)[.、]\s+(.+)$", normalized)
                if ordered_match is None:
                    continue
                ordered_item_no = int(ordered_match.group(1))
            return _HeadingMatch(
                pattern_name=pattern_name,
                line_no=0,
                explicit_level=pattern.explicit_level,
                ordered_item_no=ordered_item_no,
            )
        return None

    @staticmethod
    def _looks_like_inline_body(line: str) -> bool:
        return bool(re.search(r"[。！？；;!?]", line))

    def _build_chunks_from_blocks(
        self, text: str, blocks: list[_TextBlock]
    ) -> list[dict]:
        chunks: list[dict] = []
        heading_stack: list[_TextBlock] = []
        pending_heading_blocks: list[_TextBlock] = []
        current_body_blocks: list[_TextBlock] = []
        chunk_no = 1
        soft_body_size = max(self.chunk_size, 1)
        hard_body_size = max(
            int(soft_body_size * 1.6),
            soft_body_size + max(self.chunk_overlap, 128),
        )

        def remove_pending_heading(block: _TextBlock) -> None:
            nonlocal pending_heading_blocks
            pending_heading_blocks = [
                pending for pending in pending_heading_blocks if pending is not block
            ]

        def append_heading_chunk(block: _TextBlock) -> None:
            nonlocal chunk_no
            chunks.append(
                {
                    "chunk_no": chunk_no,
                    "chunk_text": block.text,
                    "char_start": block.start_char,
                    "char_end": block.end_char,
                    "start_line": block.start_line,
                    "end_line": block.end_line,
                }
            )
            chunk_no += 1
            remove_pending_heading(block)

        def flush_current() -> None:
            nonlocal chunk_no, current_body_blocks
            if not current_body_blocks:
                return
            body_start = current_body_blocks[0].start_char
            body_end = current_body_blocks[-1].end_char
            body_text = self._compose_body_text(current_body_blocks)
            heading_context = "\n".join(block.text for block in heading_stack)
            chunk_text = (
                f"{heading_context}\n\n{body_text}" if heading_context else body_text
            )
            chunks.append(
                {
                    "chunk_no": chunk_no,
                    "chunk_text": chunk_text,
                    "char_start": body_start,
                    "char_end": body_end,
                    "start_line": current_body_blocks[0].start_line,
                    "end_line": current_body_blocks[-1].end_line,
                }
            )
            chunk_no += 1
            for heading in heading_stack:
                remove_pending_heading(heading)
            current_body_blocks = []

        for block in blocks:
            if block.kind == "heading":
                flush_current()
                while heading_stack and (heading_stack[-1].level or 99) >= (
                    block.level or 99
                ):
                    popped_heading = heading_stack.pop()
                    if popped_heading in pending_heading_blocks:
                        append_heading_chunk(popped_heading)
                heading_stack.append(block)
                pending_heading_blocks.append(block)
                continue

            paragraph_parts = self._split_oversized_block(block, text, hard_body_size)
            for part in paragraph_parts:
                if not current_body_blocks:
                    current_body_blocks.append(part)
                    continue
                current_body_text = self._compose_body_text(current_body_blocks)
                candidate_body_text = "\n".join([current_body_text, part.text])
                candidate_size = len(candidate_body_text)
                current_size = len(current_body_text)
                should_merge = False
                if candidate_size <= soft_body_size:
                    should_merge = True
                elif (
                    candidate_size <= hard_body_size
                    and current_size < soft_body_size // 2
                ):
                    should_merge = True
                elif (
                    candidate_size <= hard_body_size
                    and len(part.text) < soft_body_size // 3
                ):
                    should_merge = True
                if should_merge:
                    current_body_blocks.append(part)
                    continue
                flush_current()
                current_body_blocks.append(part)

        flush_current()
        for heading in list(pending_heading_blocks):
            append_heading_chunk(heading)
        return chunks

    @staticmethod
    def _compose_body_text(blocks: list[_TextBlock]) -> str:
        return "\n".join(block.text for block in blocks)

    def _split_oversized_block(
        self, block: _TextBlock, text: str, max_body_size: int
    ) -> list[_TextBlock]:
        if len(block.text) <= max_body_size:
            return [block]

        table_parts = self._split_table_block(block, max_body_size)
        if table_parts is not None:
            return table_parts

        sentence_parts = self._split_block_on_sentences(block, text)
        parts: list[_TextBlock] = []
        current_group: list[_TextBlock] = []
        for sentence in sentence_parts:
            if not current_group:
                current_group.append(sentence)
                continue
            if sentence.end_char - current_group[0].start_char <= max_body_size:
                current_group.append(sentence)
                continue
            parts.append(self._merge_blocks(current_group, text))
            current_group = [sentence]
        if current_group:
            parts.append(self._merge_blocks(current_group, text))
        return parts

    def _split_table_block(
        self, block: _TextBlock, max_body_size: int
    ) -> list[_TextBlock] | None:
        """Split a Markdown table by rows, repeating header in each chunk.

        Handles both multi-line tables and single-line inline tables where
        header, separator, and data rows are concatenated with ' | ' on one line.
        """
        lines = block.text.split("\n")

        inline_result = self._try_split_inline_table(block, lines, max_body_size)
        if inline_result is not None:
            return inline_result

        if len(lines) < 3 or not lines[0].lstrip().startswith("|"):
            return None

        header_lines: list[str] = [lines[0]]
        data_start = 1
        if len(lines) > 1 and re.match(r"^\s*\|[\s:|-]+\|\s*$", lines[1]):
            header_lines.append(lines[1])
            data_start = 2

        header_text = "\n".join(header_lines)
        if len(header_text) >= max_body_size:
            return None

        parts: list[_TextBlock] = []
        current_lines: list[str] = list(header_lines)
        current_size = len(header_text)

        for i in range(data_start, len(lines)):
            line = lines[i]
            added_size = len(line) + 1
            if current_size + added_size > max_body_size and len(current_lines) > len(
                header_lines
            ):
                part_text = "\n".join(current_lines)
                parts.append(
                    _TextBlock(
                        text=part_text,
                        start_char=block.start_char,
                        end_char=block.start_char + len(part_text),
                        start_line=block.start_line,
                        end_line=block.start_line + len(current_lines) - 1,
                        kind="paragraph",
                    )
                )
                current_lines = list(header_lines)
                current_size = len(header_text)
            current_lines.append(line)
            current_size += added_size

        if len(current_lines) > len(header_lines):
            part_text = "\n".join(current_lines)
            parts.append(
                _TextBlock(
                    text=part_text,
                    start_char=block.start_char,
                    end_char=block.start_char + len(part_text),
                    start_line=block.start_line,
                    end_line=block.start_line + len(current_lines) - 1,
                    kind="paragraph",
                )
            )

        return parts if parts else None

    def _try_split_inline_table(
        self, block: _TextBlock, lines: list[str], max_body_size: int
    ) -> list[_TextBlock] | None:
        """Split single-line inline tables (header+sep+rows all on one line).

        Detects lines where an entire table (header, separator, data) is
        concatenated into a single line with pipe delimiters.
        """
        for line in lines:
            if len(line) <= max_body_size:
                continue
            if not line.lstrip().startswith("|"):
                continue
            sep_match = re.search(r"\|\s*---\s*\|", line)
            if not sep_match:
                continue

            cells = [c for c in line.split("|")]
            sep_indices = [
                i for i, c in enumerate(cells) if re.match(r"^\s*-{3,}\s*$", c)
            ]
            if len(sep_indices) < 2:
                continue

            col_count = len(sep_indices)
            first_sep_idx = sep_indices[0]
            header_cells = cells[1:first_sep_idx]
            after_sep_cells = cells[first_sep_idx + col_count :]

            header_row = "| " + " | ".join(c.strip() for c in header_cells) + " |"
            sep_row = "| " + " | ".join(["---"] * col_count) + " |"
            header_text = header_row + "\n" + sep_row

            if len(header_text) >= max_body_size:
                continue

            data_rows: list[str] = []
            for i in range(0, len(after_sep_cells) - col_count + 1, col_count):
                row_cells = after_sep_cells[i : i + col_count]
                data_rows.append("| " + " | ".join(c.strip() for c in row_cells) + " |")

            if not data_rows:
                continue

            parts: list[_TextBlock] = []
            current_chunk = header_text
            for data_row in data_rows:
                added = len(data_row) + 1
                if (
                    len(current_chunk) + added > max_body_size
                    and current_chunk != header_text
                ):
                    parts.append(
                        _TextBlock(
                            text=current_chunk,
                            start_char=block.start_char,
                            end_char=block.start_char + len(current_chunk),
                            start_line=block.start_line,
                            end_line=block.start_line,
                            kind="paragraph",
                        )
                    )
                    current_chunk = header_text
                current_chunk += "\n" + data_row

            if current_chunk != header_text:
                parts.append(
                    _TextBlock(
                        text=current_chunk,
                        start_char=block.start_char,
                        end_char=block.start_char + len(current_chunk),
                        start_line=block.start_line,
                        end_line=block.start_line,
                        kind="paragraph",
                    )
                )

            if parts:
                other_lines = [ln for ln in lines if ln is not line]
                if other_lines:
                    other_text = "\n".join(other_lines)
                    if len(other_text.strip()) > 0:
                        parts.insert(
                            0,
                            _TextBlock(
                                text=other_text,
                                start_char=block.start_char,
                                end_char=block.start_char + len(other_text),
                                start_line=block.start_line,
                                end_line=block.start_line + len(other_lines) - 1,
                                kind="paragraph",
                            ),
                        )
                return parts
        return None

    def _split_block_on_sentences(
        self, block: _TextBlock, text: str
    ) -> list[_TextBlock]:
        sentence_breaks = "。！？；;!?"
        pieces: list[_TextBlock] = []
        piece_start = block.start_char
        piece_start_line = block.start_line
        line_no = block.start_line
        cursor = block.start_char
        while cursor < block.end_char:
            char = text[cursor]
            if char in sentence_breaks:
                piece_end = cursor + 1
                piece_text = text[piece_start:piece_end]
                if piece_text.strip():
                    pieces.append(
                        _TextBlock(
                            text=piece_text,
                            start_char=piece_start,
                            end_char=piece_end,
                            start_line=piece_start_line,
                            end_line=line_no,
                            kind="paragraph",
                        )
                    )
                piece_start = piece_end
                piece_start_line = line_no + 1 if char == "\n" else line_no
            if char == "\n":
                line_no += 1
            cursor += 1

        if piece_start < block.end_char:
            pieces.append(
                _TextBlock(
                    text=text[piece_start : block.end_char],
                    start_char=piece_start,
                    end_char=block.end_char,
                    start_line=piece_start_line,
                    end_line=block.end_line,
                    kind="paragraph",
                )
            )

        if pieces:
            return [piece for piece in pieces if piece.text.strip()]

        return self._split_block_hard(block, text, max(self.chunk_size, 1))

    def _split_block_hard(
        self, block: _TextBlock, text: str, max_body_size: int
    ) -> list[_TextBlock]:
        parts: list[_TextBlock] = []
        start_char = block.start_char
        start_line = block.start_line
        line_no = block.start_line
        for cursor in range(block.start_char, block.end_char):
            if text[cursor] == "\n":
                line_no += 1
            if cursor - start_char + 1 < max_body_size and cursor + 1 < block.end_char:
                continue
            part_end = cursor + 1
            part_text = text[start_char:part_end]
            if part_text.strip():
                parts.append(
                    _TextBlock(
                        text=part_text,
                        start_char=start_char,
                        end_char=part_end,
                        start_line=start_line,
                        end_line=line_no,
                        kind="paragraph",
                    )
                )
            start_char = part_end
            start_line = line_no
        if start_char < block.end_char:
            parts.append(
                _TextBlock(
                    text=text[start_char : block.end_char],
                    start_char=start_char,
                    end_char=block.end_char,
                    start_line=start_line,
                    end_line=block.end_line,
                    kind="paragraph",
                )
            )
        return [part for part in parts if part.text.strip()]

    @staticmethod
    def _merge_blocks(blocks: list[_TextBlock], text: str) -> _TextBlock:
        return _TextBlock(
            text=text[blocks[0].start_char : blocks[-1].end_char],
            start_char=blocks[0].start_char,
            end_char=blocks[-1].end_char,
            start_line=blocks[0].start_line,
            end_line=blocks[-1].end_line,
            kind=blocks[0].kind,
            level=blocks[0].level,
        )

    def _batch_embed(self, texts: list[str]) -> list[list[float]]:
        if not texts:
            return []
        if not self.embedding_base_url:
            raise KnowledgeConfigurationError(
                "EMBEDDING_BASE_URL is required for server-side embedding"
            )
        if self.embedding_batch_max_texts == -1:
            return self._request_embeddings(texts)
        if self.embedding_batch_max_texts <= 0:
            raise KnowledgeConfigurationError(
                "embedding batch max texts must be greater than 0 or -1"
            )

        embeddings: list[list[float]] = []
        for batch_start in range(0, len(texts), self.embedding_batch_max_texts):
            embeddings.extend(
                self._request_embeddings(
                    texts[batch_start : batch_start + self.embedding_batch_max_texts]
                )
            )
        return embeddings

    def _request_embeddings(self, texts: list[str]) -> list[list[float]]:
        """Request one embedding batch from the configured endpoint."""
        headers = {"Content-Type": "application/json"}
        if self.embedding_api_key:
            headers["Authorization"] = f"Bearer {self.embedding_api_key}"

        try:
            response = httpx.post(
                f"{self.embedding_base_url.rstrip('/')}/embeddings",
                headers=headers,
                json={
                    "model": self.embedding_model_name,
                    "input": texts,
                    "dimensions": self.embedding_dimension,
                },
                timeout=self.embedding_timeout,
            )
            response.raise_for_status()
        except httpx.HTTPError as exc:
            raise KnowledgeConfigurationError(
                f"embedding service request failed: {exc}"
            ) from exc
        payload = response.json()
        data = payload.get("data") or []
        if len(data) != len(texts):
            raise KnowledgeConfigurationError(
                f"embedding API returned {len(data)} vectors for {len(texts)} inputs"
            )

        sorted_data = sorted(data, key=lambda d: d.get("index", 0))
        embeddings = [item["embedding"] for item in sorted_data]
        for index, embedding in enumerate(embeddings):
            if len(embedding) != self.embedding_dimension:
                raise KnowledgeConfigurationError(
                    f"embedding dimension mismatch at chunk {index}: "
                    f"got {len(embedding)}, expected {self.embedding_dimension}"
                )
        return embeddings
